"""
TrainHyp AI — PowerPoint Presentation Generator
Style: Clinical Minimal Light
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Design Tokens ────────────────────────────────────────────────────────────
BG       = RGBColor(0xF8, 0xFA, 0xFC)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT     = RGBColor(0x0F, 0x17, 0x2A)
MUTED    = RGBColor(0x47, 0x55, 0x69)
PRIMARY  = RGBColor(0x0B, 0x5F, 0xFF)
SUCCESS  = RGBColor(0x16, 0xA3, 0x4A)
WARNING  = RGBColor(0xD9, 0x77, 0x06)
DANGER   = RGBColor(0xDC, 0x26, 0x26)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCREENSHOT_DIR = r"C:\Users\admin\.gemini\antigravity\brain\5519c019-f23c-4a97-a63e-a8bdc8ff2e28"

SCREENSHOTS = {
    "overview":  os.path.join(SCREENSHOT_DIR, "data_overview_screenshot_1777259812712.png"),
    "volume":    os.path.join(SCREENSHOT_DIR, "volume_hypertrophy_page_1777259855097.png"),
    "optimizer": os.path.join(SCREENSHOT_DIR, "optimizer_custom_profile_1777260008665.png"),
    "casestudy": os.path.join(SCREENSHOT_DIR, "case_study_report_1777260144392.png"),
}


def set_slide_bg(slide, color=BG):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Inter"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Subtle shadow
    return shape


def add_accent_line(slide, left, top, width, color=PRIMARY):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_stat_card(slide, left, top, value, label, accent_color=PRIMARY):
    """Add a stat card (value + label)."""
    card_w = Inches(2.6)
    card_h = Inches(1.35)
    rect = add_rounded_rect(slide, left, top, card_w, card_h, SURFACE)

    # Accent line at top
    add_accent_line(slide, left + Inches(0.2), top + Inches(0.12), Inches(0.5), accent_color)

    # Value
    add_text_box(slide, left + Inches(0.2), top + Inches(0.3), card_w - Inches(0.4), Inches(0.6),
                 str(value), font_size=28, bold=True, color=accent_color)

    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(0.85), card_w - Inches(0.4), Inches(0.4),
                 label, font_size=11, color=MUTED)


def add_bullet_list(slide, left, top, width, items, font_size=14, color=TEXT, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Use bullet character
        if isinstance(item, tuple):
            # (bullet_char, text)
            p.text = f"{item[0]}  {item[1]}"
        else:
            p.text = f"•  {item}"

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.space_after = spacing
        p.space_before = Pt(2)

    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, SURFACE)

    # Left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRIMARY
    strip.line.fill.background()

    # Top decorative bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), 0, SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # Course badge
    badge = add_rounded_rect(slide, Inches(1.2), Inches(1.0), Inches(4.0), Inches(0.45), LIGHT_BLUE)
    add_text_box(slide, Inches(1.4), Inches(1.05), Inches(3.6), Inches(0.4),
                 "STAT3013 — STATISTICAL ANALYSIS", font_size=12, bold=True, color=PRIMARY)

    # Title
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
                 "TrainHyp AI", font_size=48, bold=True, color=TEXT)

    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(0.8),
                 "Training Volume Optimizer for Muscle Hypertrophy",
                 font_size=24, color=MUTED)

    # Topic line
    add_text_box(slide, Inches(1.2), Inches(3.6), Inches(11), Inches(0.6),
                 "Statistical Analysis of Training Volume and Its Correlation with Muscle Hypertrophy",
                 font_size=14, color=MUTED)

    # Divider
    add_accent_line(slide, Inches(1.2), Inches(4.4), Inches(2.0), PRIMARY)

    # Team info
    team_lines = [
        "Group 7  ·  STAT3013.Q21.CTTT",
        "",
        "24522046  Võ Tấn Vũ          24521521  Nguyễn Đình Sang",
        "24521446  Nguyễn Phan Hoàng Quân    24521450  Phạm Như Quân",
        "",
        "Lecturers:  Dr. Trần Văn Hải Triều  &  TA. Nguyễn Minh Nhựt",
    ]
    y = Inches(4.7)
    for line in team_lines:
        if line == "":
            y += Inches(0.15)
            continue
        sz = 11 if "Lecturers" in line else 12
        clr = MUTED if "Lecturers" in line else TEXT
        add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.35),
                     line, font_size=sz, color=clr)
        y += Inches(0.32)

    # Year
    add_text_box(slide, Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4),
                 "2026", font_size=14, bold=True, color=MUTED, alignment=PP_ALIGN.RIGHT)


def slide_02_problem(prs):
    """The Problem slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "01  INTRODUCTION", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "The Problem", font_size=36, bold=True, color=TEXT)

    # Key question
    q_rect = add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0), LIGHT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(11.2), Inches(0.8),
                 "\"Given an individual's training profile, what is the optimal weekly volume "
                 "(sets/week) to maximize hypertrophic response?\"",
                 font_size=16, bold=False, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # Two columns
    # Left: current approach
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.5),
                 "❌  Current Approach", font_size=16, bold=True, color=DANGER)

    items_left = [
        "Population-level guidelines: \"10–20 sets/week\"",
        "One-size-fits-all — ignores individual variation",
        "No uncertainty quantification",
        "No safety guardrails for edge cases",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(3.8), Inches(5.2), items_left, font_size=13, color=MUTED)

    # Right: our approach
    add_text_box(slide, Inches(7.0), Inches(3.2), Inches(5.5), Inches(0.5),
                 "✅  TrainHyp AI Approach", font_size=16, bold=True, color=SUCCESS)

    items_right = [
        "Personalized prediction via 4-model ensemble",
        "13-feature clinical profile → Hedges' g",
        "Probabilistic uncertainty (NGBoost σ + 95% CI)",
        "4-rule safety engine for OOD detection",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(3.8), Inches(5.5), items_right, font_size=13, color=MUTED)

    # Bottom stat
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
                 "Data Source:  69 RCTs  ·  199 muscle-level observations  ·  measured via ultrasound & MRI",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


def slide_03_objectives(prs):
    """Research Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "01  INTRODUCTION", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Research Objectives", font_size=36, bold=True, color=TEXT)

    objectives = [
        ("01", "Analyse dose–response relationships from 69 RCTs (199 observations) for muscle hypertrophy"),
        ("02", "Build a 4-model ML ensemble (EBM · NGBoost · CatBoost · GPR) to predict Hedges' g effect sizes"),
        ("03", "Create an interactive clinical decision-support tool with personalized volume recommendations"),
        ("04", "Implement probabilistic uncertainty estimates and safety guardrails for out-of-distribution inputs"),
    ]

    y = Inches(2.0)
    for num, text in objectives:
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = PRIMARY
        circ.line.fill.background()
        add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(0.55), Inches(0.4),
                     num, font_size=16, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)

        # Text
        add_text_box(slide, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.5),
                     text, font_size=15, color=TEXT)

        y += Inches(1.15)


def slide_04_dataset(prs):
    """Dataset Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "02  DATA & METHODOLOGY", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Dataset Overview", font_size=36, bold=True, color=TEXT)

    # Stat cards
    add_stat_card(slide, Inches(0.8), Inches(2.0), "69", "RCT Studies\n(1989–2024)", PRIMARY)
    add_stat_card(slide, Inches(3.7), Inches(2.0), "199", "Muscle-level\nObservations", PRIMARY)
    add_stat_card(slide, Inches(6.6), Inches(2.0), "13", "Input Features", PRIMARY)
    add_stat_card(slide, Inches(9.5), Inches(2.0), "Hedges' g", "Target Variable\n(Effect Size)", SUCCESS)

    # Feature categories
    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
                 "Feature Groups", font_size=18, bold=True, color=TEXT)

    # 3 feature group cards
    groups = [
        ("Training Variables (7)", "sets/week (all), sets/week (direct),\nfrequency, sessions, rep range,\ninter-set rest, % sets to failure", PRIMARY),
        ("Subject Profile (3)", "age, sex, training status\n(untrained vs trained)", WARNING),
        ("Context Variables (3)", "upper/lower body, nutrition control,\nprogram duration (weeks)", SUCCESS),
    ]

    x = Inches(0.8)
    for title, desc, color in groups:
        card = add_rounded_rect(slide, x, Inches(4.4), Inches(3.7), Inches(2.2), SURFACE)
        add_accent_line(slide, x + Inches(0.15), Inches(4.55), Inches(0.8), color)
        add_text_box(slide, x + Inches(0.2), Inches(4.7), Inches(3.3), Inches(0.4),
                     title, font_size=13, bold=True, color=TEXT)
        add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(3.3), Inches(1.5),
                     desc, font_size=11, color=MUTED)
        x += Inches(4.0)


def slide_05_hedges_g(prs):
    """Hedges' g explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "02  DATA & METHODOLOGY", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Why Hedges' g?", font_size=36, bold=True, color=TEXT)

    # Formula box
    formula_rect = add_rounded_rect(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4), SURFACE)
    add_text_box(slide, Inches(1.2), Inches(2.1), Inches(11), Inches(0.5),
                 "Hedges' g  =  J  ×  (post.mean − pre.mean) / pooled SD",
                 font_size=18, bold=True, color=TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(2.7), Inches(11), Inches(0.6),
                 "Bias-corrected standardized mean difference  —  suitable for small-sample studies\n"
                 "Standardizes across measurement methods (ultrasound, MRI, biopsy)",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Classification table
    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.5),
                 "Responder Classification", font_size=18, bold=True, color=TEXT)

    classes = [
        ("Low Responder", "g < 0.2", "Negligible to small effect", DANGER),
        ("Medium Responder", "0.2 ≤ g < 0.8", "Moderate hypertrophic effect", WARNING),
        ("High Responder", "g ≥ 0.8", "Large hypertrophic effect", SUCCESS),
    ]

    y = Inches(4.4)
    for name, threshold, desc, color in classes:
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.75), SURFACE)
        # Color indicator
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), y + Inches(0.1), Inches(0.08), Inches(0.55))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()

        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     name, font_size=14, bold=True, color=TEXT)
        add_text_box(slide, Inches(4.0), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     threshold, font_size=14, bold=True, color=color)
        add_text_box(slide, Inches(7.0), y + Inches(0.1), Inches(5), Inches(0.4),
                     desc, font_size=13, color=MUTED)
        y += Inches(0.9)


def slide_06_ensemble(prs):
    """4-Model Ensemble pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "02  DATA & METHODOLOGY", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "4-Model Ensemble Architecture", font_size=36, bold=True, color=TEXT)

    # Input box
    input_box = add_rounded_rect(slide, Inches(0.5), Inches(2.2), Inches(2.2), Inches(1.0), PRIMARY)
    add_text_box(slide, Inches(0.6), Inches(2.3), Inches(2.0), Inches(0.4),
                 "INPUT", font_size=10, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), Inches(2.6), Inches(2.0), Inches(0.5),
                 "13 Features\n+ Preprocessing", font_size=13, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)

    # Arrow
    add_text_box(slide, Inches(2.8), Inches(2.5), Inches(0.5), Inches(0.4),
                 "→", font_size=28, bold=True, color=MUTED, alignment=PP_ALIGN.CENTER)

    # 4 model boxes
    models = [
        ("EBM-GAM", "Regression", "Hedges' g\nprediction", PRIMARY),
        ("NGBoost", "Uncertainty", "σ + 95% CI", WARNING),
        ("CatBoost", "Classification", "High / Med / Low\nresponder", SUCCESS),
        ("GPR", "OOD Detection", "σ_gpr\n(input novelty)", DANGER),
    ]

    y = Inches(1.8)
    for name, task, output, color in models:
        card = add_rounded_rect(slide, Inches(3.5), y, Inches(3.2), Inches(1.1), SURFACE)
        # Color left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), y, Inches(0.08), Inches(1.1))
        border.fill.solid()
        border.fill.fore_color.rgb = color
        border.line.fill.background()

        add_text_box(slide, Inches(3.8), y + Inches(0.08), Inches(1.5), Inches(0.35),
                     name, font_size=14, bold=True, color=color)
        add_text_box(slide, Inches(5.3), y + Inches(0.08), Inches(1.2), Inches(0.35),
                     task, font_size=10, color=MUTED)
        add_text_box(slide, Inches(3.8), y + Inches(0.5), Inches(2.7), Inches(0.6),
                     output, font_size=11, color=TEXT)
        y += Inches(1.25)

    # Arrow to output
    add_text_box(slide, Inches(6.8), Inches(3.2), Inches(0.5), Inches(0.4),
                 "→", font_size=28, bold=True, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Rule engine + output box
    rule_box = add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(2.3), Inches(1.8), SURFACE)
    add_accent_line(slide, Inches(7.6), Inches(1.95), Inches(0.6), PRIMARY)
    add_text_box(slide, Inches(7.6), Inches(2.1), Inches(2.1), Inches(0.35),
                 "Rule Engine", font_size=14, bold=True, color=TEXT)
    items = ["Range check", "σ_ngb threshold", "σ_gpr threshold", "Volume bounds"]
    add_bullet_list(slide, Inches(7.6), Inches(2.5), Inches(2.0), items, font_size=10, color=MUTED, spacing=Pt(2))

    # Arrow
    add_text_box(slide, Inches(9.9), Inches(2.5), Inches(0.5), Inches(0.4),
                 "→", font_size=28, bold=True, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Final output
    out_box = add_rounded_rect(slide, Inches(10.5), Inches(1.8), Inches(2.5), Inches(1.8), PRIMARY)
    add_text_box(slide, Inches(10.6), Inches(1.9), Inches(2.3), Inches(0.35),
                 "OUTPUT", font_size=10, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)
    outputs = [
        "• Hedges' g score",
        "• Responder class",
        "• Optimal sets/week",
        "• 95% CI",
        "• Safety warnings",
    ]
    yy = Inches(2.3)
    for o in outputs:
        add_text_box(slide, Inches(10.7), yy, Inches(2.2), Inches(0.3),
                     o, font_size=11, color=SURFACE)
        yy += Inches(0.26)

    # Dose response sweep
    sweep_box = add_rounded_rect(slide, Inches(7.5), Inches(4.0), Inches(5.5), Inches(1.2), LIGHT_BLUE)
    add_text_box(slide, Inches(7.7), Inches(4.1), Inches(5.1), Inches(0.4),
                 "Dose–Response Sweep", font_size=14, bold=True, color=PRIMARY)
    add_text_box(slide, Inches(7.7), Inches(4.5), Inches(5.1), Inches(0.6),
                 "EBM prediction across 1–50 sets/week  →  find argmax (P90-capped ≤ 32 sets)\n"
                 "Generates personalized dose–response curve for visualization",
                 font_size=11, color=MUTED)

    # Bottom note
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                 "Each model contributes a unique signal: EBM predicts magnitude · NGBoost quantifies spread · "
                 "CatBoost classifies response · GPR detects novelty",
                 font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)


def slide_07_safety(prs):
    """Safety Guardrails."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "02  DATA & METHODOLOGY", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Safety Guardrails — 4-Rule Engine", font_size=36, bold=True, color=TEXT)

    rules = [
        ("⚠️", "Feature Range Extrapolation",
         "Each input checked against training data min/max.\n"
         "Flags when any feature falls outside observed range.",
         DANGER),
        ("📊", "NGBoost High Uncertainty",
         "If predictive σ exceeds threshold (P75/P90 of training σ),\n"
         "the point prediction is flagged as unreliable.",
         WARNING),
        ("🔍", "GPR Out-of-Distribution",
         "GPR posterior variance grows with distance from training data.\n"
         "Measures input novelty — complements NGBoost uncertainty.",
         PRIMARY),
        ("🏋️", "Volume Bounds",
         "Too low → junk volume warning. Too high → overtraining risk.\n"
         "Based on physiological boundaries in meta.pkl.",
         MUTED),
    ]

    y = Inches(2.0)
    for icon, title, desc, color in rules:
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), SURFACE)

        # Color indicator
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), y + Inches(0.1), Inches(0.08), Inches(0.95))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()

        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(0.5), Inches(0.5),
                     icon, font_size=20, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.8), y + Inches(0.12), Inches(4), Inches(0.4),
                     title, font_size=15, bold=True, color=TEXT)
        add_text_box(slide, Inches(1.8), y + Inches(0.5), Inches(10), Inches(0.7),
                     desc, font_size=12, color=MUTED)

        y += Inches(1.3)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
                 "Warnings are non-blocking — predictions always returned. "
                 "Guardrails inform, not restrict.",
                 font_size=12, bold=True, color=MUTED, alignment=PP_ALIGN.CENTER)


def slide_08_architecture(prs):
    """System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "03  SYSTEM ARCHITECTURE", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "System Architecture", font_size=36, bold=True, color=TEXT)

    # Frontend box
    fe_box = add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.5), SURFACE)
    add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5), Inches(0.4),
                 "FRONTEND  —  React 19 + TypeScript", font_size=13, bold=True, color=PRIMARY)

    pages = [
        ("📊", "Data Overview", "Histogram + Boxplot + Filters"),
        ("📈", "Volume vs Hypertrophy", "Scatter + Pearson r"),
        ("🎯", "Volume Optimizer", "AI Prediction + Dose-Response"),
        ("🔬", "Case Study", "Clinical Report Generation"),
    ]

    py = Inches(2.7)
    for icon, name, desc in pages:
        add_text_box(slide, Inches(0.9), py, Inches(0.5), Inches(0.4), icon, font_size=16)
        add_text_box(slide, Inches(1.4), py, Inches(2), Inches(0.35), name, font_size=12, bold=True, color=TEXT)
        add_text_box(slide, Inches(3.5), py, Inches(2.3), Inches(0.35), desc, font_size=10, color=MUTED)
        py += Inches(0.55)

    # Key patterns
    add_text_box(slide, Inches(0.9), Inches(5.2), Inches(5), Inches(0.4),
                 "React Router · Recharts · Motion · Tailwind v4",
                 font_size=10, color=MUTED)
    add_text_box(slide, Inches(0.9), Inches(5.55), Inches(5), Inches(0.4),
                 "ErrorBoundary · Graceful Degradation · A11y",
                 font_size=10, color=MUTED)

    # Arrow
    add_text_box(slide, Inches(6.1), Inches(3.7), Inches(0.8), Inches(0.6),
                 "HTTP\nJSON", font_size=11, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.0), Inches(3.2), Inches(1.0), Inches(0.5),
                 "⟷", font_size=32, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # Backend box
    be_box = add_rounded_rect(slide, Inches(7.2), Inches(2.0), Inches(5.8), Inches(4.5), SURFACE)
    add_text_box(slide, Inches(7.4), Inches(2.1), Inches(5.4), Inches(0.4),
                 "BACKEND  —  FastAPI + Python", font_size=13, bold=True, color=PRIMARY)

    # API endpoints
    add_text_box(slide, Inches(7.4), Inches(2.6), Inches(5.4), Inches(0.35),
                 "API Endpoints", font_size=12, bold=True, color=TEXT)

    endpoints = [
        ("GET", "/health", "Model status"),
        ("GET", "/model-info", "SHAP importance"),
        ("POST", "/predict", "Full prediction"),
    ]
    ey = Inches(3.0)
    for method, path, desc in endpoints:
        clr = SUCCESS if method == "GET" else WARNING
        add_text_box(slide, Inches(7.6), ey, Inches(0.7), Inches(0.3), method, font_size=9, bold=True, color=clr)
        add_text_box(slide, Inches(8.3), ey, Inches(1.8), Inches(0.3), path, font_size=10, bold=True, color=TEXT)
        add_text_box(slide, Inches(10.2), ey, Inches(2.5), Inches(0.3), desc, font_size=10, color=MUTED)
        ey += Inches(0.35)

    # Model store
    add_text_box(slide, Inches(7.4), Inches(4.3), Inches(5.4), Inches(0.35),
                 "AI Engine  →  4 Models  →  28 .pkl files", font_size=12, bold=True, color=TEXT)
    add_text_box(slide, Inches(7.4), Inches(4.7), Inches(5.4), Inches(0.35),
                 "EBM · NGBoost · CatBoost · GPR", font_size=11, color=MUTED)
    add_text_box(slide, Inches(7.4), Inches(5.1), Inches(5.4), Inches(0.35),
                 "Pydantic v2 validation · CORS · Uvicorn",
                 font_size=10, color=MUTED)
    add_text_box(slide, Inches(7.4), Inches(5.45), Inches(5.4), Inches(0.35),
                 "Auto Swagger docs at /docs",
                 font_size=10, color=MUTED)


def slide_09_techstack(prs):
    """Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "03  SYSTEM ARCHITECTURE", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Technology Stack", font_size=36, bold=True, color=TEXT)

    categories = [
        ("Frontend", [
            ("React 19", "UI component framework"),
            ("TypeScript (strict)", "Static type safety"),
            ("Vite 6", "Build tool + HMR"),
            ("Tailwind CSS v4", "Utility-first CSS"),
            ("Recharts 3", "Declarative charts"),
            ("React Router v7", "URL-based navigation"),
        ], PRIMARY),
        ("Backend", [
            ("FastAPI", "Async web framework"),
            ("Pydantic v2", "Request validation"),
            ("Uvicorn", "ASGI server"),
            ("Pandas / NumPy", "Data operations"),
            ("scikit-learn", "Scaler + GPR + Imputer"),
            ("Joblib", "Model serialization"),
        ], SUCCESS),
        ("ML Models", [
            ("EBM (interpret)", "Regression — Hedges' g"),
            ("NGBoost", "Probabilistic σ + CI"),
            ("CatBoost", "Multi-class classifier"),
            ("GPR", "OOD detection"),
            ("SHAP", "Feature importance"),
            ("", ""),
        ], WARNING),
    ]

    x = Inches(0.5)
    for cat_name, items, color in categories:
        card = add_rounded_rect(slide, x, Inches(2.0), Inches(3.9), Inches(5.0), SURFACE)
        add_accent_line(slide, x + Inches(0.15), Inches(2.15), Inches(0.8), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.4),
                     cat_name, font_size=16, bold=True, color=color)

        iy = Inches(2.9)
        for tech, desc in items:
            if not tech:
                continue
            add_text_box(slide, x + Inches(0.2), iy, Inches(2.0), Inches(0.3),
                         tech, font_size=12, bold=True, color=TEXT)
            add_text_box(slide, x + Inches(0.2), iy + Inches(0.25), Inches(3.5), Inches(0.3),
                         desc, font_size=10, color=MUTED)
            iy += Inches(0.55)

        x += Inches(4.15)


def _add_screenshot_slide(prs, section_num, section_label, title, subtitle, img_path):
    """Generic screenshot slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.3),
                 f"{section_num}  {section_label}", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                 title, font_size=24, bold=True, color=TEXT)

    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                     subtitle, font_size=12, color=MUTED)

    # Screenshot with border
    if os.path.exists(img_path):
        img_top = Inches(1.5)
        img_height = Inches(5.6)
        img_width = Inches(10.5)

        # Shadow rectangle behind image
        shadow = add_rounded_rect(slide, Inches(1.35), img_top + Inches(0.05),
                                  img_width + Inches(0.1), img_height + Inches(0.1),
                                  RGBColor(0xE2, 0xE8, 0xF0))

        slide.shapes.add_picture(img_path, Inches(1.4), img_top, img_width, img_height)


def slide_10_demo_overview(prs):
    _add_screenshot_slide(prs, "04", "APPLICATION DEMO", "Data Overview",
                          "Interactive dataset explorer — 69 RCTs · 198 observations · histogram + boxplot + filters",
                          SCREENSHOTS["overview"])


def slide_11_demo_volume(prs):
    _add_screenshot_slide(prs, "04", "APPLICATION DEMO", "Volume vs Hypertrophy",
                          "Scatter analysis — 198 observations · Pearson r = +0.44 · class-colored · Cohen's benchmarks",
                          SCREENSHOTS["volume"])


def slide_12_demo_optimizer(prs):
    _add_screenshot_slide(prs, "04", "APPLICATION DEMO", "Volume Optimizer — Custom Profile",
                          "10 adjustable parameters · real-time AI prediction · dose-response curve · SHAP importance",
                          SCREENSHOTS["optimizer"])


def slide_13_demo_casestudy(prs):
    _add_screenshot_slide(prs, "04", "APPLICATION DEMO", "Case Study — Clinical Report",
                          "Pre-defined clinical cases · full AI analysis · dose-response curve · clinical recommendations",
                          SCREENSHOTS["casestudy"])


def slide_14_key_findings(prs):
    """Key Findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "05  RESULTS", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Key Findings", font_size=36, bold=True, color=TEXT)

    findings = [
        ("01", "Sets/week is the strongest predictor",
         "SHAP importance = 0.42 — dominates all other features.\n"
         "Pearson r = +0.44 between sets/week and Hedges' g.",
         PRIMARY),
        ("02", "Optimal range: 20–28 sets/week for most profiles",
         "EBM dose-response curve peaks in this range for trained individuals.\n"
         "Results consistent with meta-analytic evidence (Schoenfeld et al.).",
         SUCCESS),
        ("03", "Diminishing returns above P90 (~32 sets/week)",
         "Only 5 observations with sets ≥ 40 in dataset.\n"
         "P90 cap prevents unreliable extrapolation beyond data support.",
         WARNING),
        ("04", "Training status is the 2nd most important factor",
         "SHAP importance = 0.28. Untrained subjects show higher effect sizes\n"
         "at lower volumes — consistent with \"newbie gains\" phenomenon.",
         DANGER),
    ]

    y = Inches(2.0)
    for num, title, desc, color in findings:
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.1), SURFACE)

        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.25), Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        add_text_box(slide, Inches(1.1), y + Inches(0.33), Inches(0.5), Inches(0.35),
                     num, font_size=14, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.9), y + Inches(0.08), Inches(10), Inches(0.35),
                     title, font_size=15, bold=True, color=TEXT)
        add_text_box(slide, Inches(1.9), y + Inches(0.45), Inches(10.3), Inches(0.65),
                     desc, font_size=11, color=MUTED)

        y += Inches(1.25)


def slide_15_validation(prs):
    """Model Validation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "05  RESULTS", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Model Validation — Test Cases", font_size=36, bold=True, color=TEXT)

    # Test case table header
    headers = ["Case", "Description", "Expected", "Predicted", "Match"]
    col_widths = [Inches(0.8), Inches(4.5), Inches(1.5), Inches(1.5), Inches(1.0)]
    col_x = [Inches(0.8)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # Header row
    header_bg = add_rounded_rect(slide, Inches(0.6), Inches(2.0), Inches(12.0), Inches(0.55), PRIMARY)
    for i, h in enumerate(headers):
        add_text_box(slide, col_x[i], Inches(2.05), col_widths[i], Inches(0.45),
                     h, font_size=12, bold=True, color=SURFACE, alignment=PP_ALIGN.CENTER)

    # Data rows
    cases = [
        ("#1", "Trained male, moderate volume, upper body", "Medium", "Medium", "✓"),
        ("#2", "High frequency, high volume trained", "High", "High", "✓"),
        ("#3", "Untrained, very low volume", "Low", "Low", "✓"),
        ("#4", "Long program, untrained, moderate volume", "Medium", "Medium", "✓"),
        ("#5", "Very high volume — potential overtraining", "High", "High", "✓"),
    ]

    y = Inches(2.65)
    for case_data in cases:
        row_bg = add_rounded_rect(slide, Inches(0.6), y, Inches(12.0), Inches(0.6), SURFACE)
        for i, val in enumerate(case_data):
            clr = SUCCESS if val == "✓" else TEXT
            bld = True if i == 0 or val == "✓" else False
            add_text_box(slide, col_x[i], y + Inches(0.08), col_widths[i], Inches(0.45),
                         val, font_size=12, bold=bld, color=clr, alignment=PP_ALIGN.CENTER)
        y += Inches(0.7)

    # Summary
    summary_box = add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), LIGHT_BLUE)
    add_text_box(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
                 "5/5 test cases correctly classified  ·  95% average confidence  ·  "
                 "Graceful fallback ensures usability when backend is offline",
                 font_size=13, bold=False, color=PRIMARY, alignment=PP_ALIGN.CENTER)


def slide_16_limitations(prs):
    """Limitations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "06  DISCUSSION", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Limitations", font_size=36, bold=True, color=TEXT)

    limitations = [
        ("Dataset Size",
         "199 observations from 69 studies — study-level aggregates,\n"
         "not individual-level data. Limited statistical power for rare subgroups."),
        ("Feature Representation",
         "Sex (sex.male) is a study-level proportion, not individual binary.\n"
         "Some features (e.g., years.exp) had significant missingness (45%)."),
        ("Cross-Sectional Only",
         "No longitudinal tracking — predictions are for a single time point.\n"
         "Cannot model progressive overload or adaptation over time."),
        ("Publication Bias",
         "Dataset comprises published RCTs only.\n"
         "Studies with negative/null results may be underrepresented."),
    ]

    y = Inches(2.0)
    for title, desc in limitations:
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.05), SURFACE)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.35),
                     f"⚠  {title}", font_size=14, bold=True, color=WARNING)
        add_text_box(slide, Inches(1.5), y + Inches(0.45), Inches(10.5), Inches(0.6),
                     desc, font_size=12, color=MUTED)
        y += Inches(1.2)


def slide_17_future(prs):
    """Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "06  DISCUSSION", font_size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Future Work", font_size=36, bold=True, color=TEXT)

    items = [
        ("📦", "Expand Dataset",
         "Incorporate newer studies and individual-level data.\n"
         "Add muscle-group-specific volume analysis."),
        ("📱", "Mobile & Accessibility",
         "Responsive mobile layout + dark mode.\n"
         "Enhanced WCAG 2.1 compliance."),
        ("☁️", "Cloud Deployment",
         "Deploy to Vercel (frontend) + Railway (backend).\n"
         "Public access without local installation."),
        ("🔬", "Advanced Models",
         "Experiment with Transformer-based architectures.\n"
         "Add longitudinal prediction capability."),
    ]

    x = Inches(0.5)
    y_row1 = Inches(2.0)
    y_row2 = Inches(4.2)

    for i, (icon, title, desc) in enumerate(items):
        cx = Inches(0.5) + (i % 2) * Inches(6.2)
        cy = y_row1 if i < 2 else y_row2

        card = add_rounded_rect(slide, cx, cy, Inches(5.8), Inches(1.8), SURFACE)
        add_text_box(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(0.5), Inches(0.5),
                     icon, font_size=24)
        add_text_box(slide, cx + Inches(0.8), cy + Inches(0.15), Inches(4.5), Inches(0.4),
                     title, font_size=16, bold=True, color=TEXT)
        add_text_box(slide, cx + Inches(0.8), cy + Inches(0.6), Inches(4.8), Inches(1.0),
                     desc, font_size=12, color=MUTED)


def slide_18_thankyou(prs):
    """Thank You / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)

    # Left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRIMARY
    strip.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), 0, SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(1.0),
                 "Thank You", font_size=52, bold=True, color=TEXT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(0.6),
                 "Questions & Discussion", font_size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Divider
    add_accent_line(slide, Inches(5.5), Inches(3.4), Inches(2.3), PRIMARY)

    # Project info
    info_box = add_rounded_rect(slide, Inches(3.5), Inches(4.0), Inches(6.3), Inches(2.5), BG)

    info_lines = [
        ("TrainHyp AI — Training Volume Optimizer", 14, True, TEXT),
        ("", 8, False, TEXT),
        ("Group 7  ·  STAT3013.Q21.CTTT", 12, False, MUTED),
        ("", 6, False, TEXT),
        ("Võ Tấn Vũ  ·  Nguyễn Đình Sang", 11, False, TEXT),
        ("Nguyễn Phan Hoàng Quân  ·  Phạm Như Quân", 11, False, TEXT),
        ("", 6, False, TEXT),
        ("Dr. Trần Văn Hải Triều  &  TA. Nguyễn Minh Nhựt", 11, False, MUTED),
    ]

    iy = Inches(4.1)
    for text, sz, bld, clr in info_lines:
        if text == "":
            iy += Inches(0.12)
            continue
        add_text_box(slide, Inches(3.7), iy, Inches(5.9), Inches(0.35),
                     text, font_size=sz, bold=bld, color=clr, alignment=PP_ALIGN.CENTER)
        iy += Inches(0.3)

    # Footer
    add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11), Inches(0.4),
                 "STAT3013 — Statistical Analysis  ·  2026",
                 font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building slides...")

    slide_01_title(prs)
    print("  ✓ Slide 01: Title")

    slide_02_problem(prs)
    print("  ✓ Slide 02: The Problem")

    slide_03_objectives(prs)
    print("  ✓ Slide 03: Research Objectives")

    slide_04_dataset(prs)
    print("  ✓ Slide 04: Dataset Overview")

    slide_05_hedges_g(prs)
    print("  ✓ Slide 05: Hedges' g Explanation")

    slide_06_ensemble(prs)
    print("  ✓ Slide 06: 4-Model Ensemble")

    slide_07_safety(prs)
    print("  ✓ Slide 07: Safety Guardrails")

    slide_08_architecture(prs)
    print("  ✓ Slide 08: System Architecture")

    slide_09_techstack(prs)
    print("  ✓ Slide 09: Tech Stack")

    slide_10_demo_overview(prs)
    print("  ✓ Slide 10: Demo — Data Overview")

    slide_11_demo_volume(prs)
    print("  ✓ Slide 11: Demo — Volume vs Hypertrophy")

    slide_12_demo_optimizer(prs)
    print("  ✓ Slide 12: Demo — Volume Optimizer")

    slide_13_demo_casestudy(prs)
    print("  ✓ Slide 13: Demo — Case Study")

    slide_14_key_findings(prs)
    print("  ✓ Slide 14: Key Findings")

    slide_15_validation(prs)
    print("  ✓ Slide 15: Model Validation")

    slide_16_limitations(prs)
    print("  ✓ Slide 16: Limitations")

    slide_17_future(prs)
    print("  ✓ Slide 17: Future Work")

    slide_18_thankyou(prs)
    print("  ✓ Slide 18: Thank You")

    output_path = r"f:\class_project\pttk\TrainHyp_Slides.pptx"
    prs.save(output_path)
    print(f"\n✅ Saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
